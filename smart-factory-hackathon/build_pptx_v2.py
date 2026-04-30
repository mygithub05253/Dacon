#!/usr/bin/env python3
"""Build QualityLens 기획서 PPTX v2 — 텍스트 다이어트 + 서식 개선."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = "/sessions/clever-nifty-volta/mnt/Dacon/smart-factory-hackathon/2026-산출물-양식.pptx"
OUTPUT = "/sessions/clever-nifty-volta/mnt/Dacon/smart-factory-hackathon/QualityLens_기획서.pptx"

# ── Colors ──
C_DARK = RGBColor(0x20, 0x21, 0x24)
C_HEADER = RGBColor(0x1A, 0x47, 0x8A)   # 진한 파랑 (헤더)
C_ACCENT = RGBColor(0x2E, 0x7D, 0x32)   # 그린 (강조 포인트)
C_BODY = RGBColor(0x33, 0x33, 0x33)      # 본문 텍스트
C_RED = RGBColor(0xC6, 0x28, 0x28)       # 경고/중요
C_MUTED = RGBColor(0x66, 0x66, 0x66)     # 보조 설명

FONT = "Malgun Gothic"  # LO에서 NanumGothic으로 자동 매핑

prs = Presentation(TEMPLATE)


def clear_and_set(text_frame, paragraphs_data, font_size=10, font_name=FONT):
    """Clear existing paragraphs and set new content."""
    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
        p = text_frame.paragraphs[i]._p
        p.getparent().remove(p)

    first = True
    for pd in paragraphs_data:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()

        p.text = pd.get("text", "")
        p.font.size = pd.get("size", Pt(font_size))
        p.font.name = font_name
        p.font.bold = pd.get("bold", False)
        p.font.color.rgb = pd.get("color", C_BODY)
        p.alignment = pd.get("align", PP_ALIGN.LEFT)

        if pd.get("indent", 0) > 0:
            p.level = pd.get("indent", 0)

        p.space_after = Pt(pd.get("space_after", 3))
        p.space_before = Pt(pd.get("space_before", 0))


def find_guide_textbox(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "작성방법" in shape.text_frame.text:
                return shape
    return None


def H(text, sa=1):
    """섹션 헤더 (14pt, 파란색, 굵게)"""
    return {"text": text, "bold": True, "size": Pt(14), "color": C_HEADER, "space_after": sa, "space_before": 4}

def SH(text, sa=1):
    """서브 헤더 (12pt, 진한색, 굵게)"""
    return {"text": text, "bold": True, "size": Pt(12), "color": C_DARK, "space_after": sa, "space_before": 3}

def B(text, sa=2):
    """본문 텍스트 (10pt)"""
    return {"text": text, "bold": False, "size": Pt(10), "color": C_BODY, "space_after": sa}

def BL(text, sa=2):
    """불렛 포인트 (10pt, • 포함)"""
    return {"text": f"• {text}", "bold": False, "size": Pt(10), "color": C_BODY, "space_after": sa}

def BLB(label, desc, sa=2):
    """불렛 + 볼드 라벨 (라벨: 설명)"""
    return {"text": f"• {label}: {desc}", "bold": False, "size": Pt(10), "color": C_BODY, "space_after": sa}

def ACC(text, sa=2):
    """강조 텍스트 (10pt, 빨간색, 굵게)"""
    return {"text": text, "bold": True, "size": Pt(10), "color": C_RED, "space_after": sa}

def SP(h=4):
    """빈 줄 (간격용)"""
    return {"text": "", "size": Pt(4), "space_after": h}


# Position constants for content textboxes
TB_TOP = Emu(1600000)
TB_LEFT = Emu(700000)
TB_WIDTH = Emu(10800000)
TB_HEIGHT = Emu(4800000)


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_table:
        table = shape.table
        table.cell(0, 1).text = "QualityLens"
        for p in table.cell(0, 1).text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = FONT
            p.font.bold = True
        table.cell(1, 1).text = "QualityLens — AI 기반 스마트 품질 예측·원인 분석 통합 플랫폼"
        for p in table.cell(1, 1).text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = FONT

print("Slide 1 (Title) - Done")


# ═══════════════════════════════════════════════
# SLIDE 2: 문제 정의 — 텍스트 다이어트 적용
# ═══════════════════════════════════════════════
slide2 = prs.slides[1]
guide2 = find_guide_textbox(slide2)
if guide2:
    clear_and_set(guide2.text_frame, [
        H("해결하고자 하는 제조 현장의 문제"),
        B("경기도 내 중소 전자부품 제조사들은 수백 개 공정 센서 데이터를 보유하면서도, 전문 분석 인력 부재로 최종 검사 이후에야 불량을 확인하는 '사후 대응'에 의존. 데이터 기반 예방적 품질 관리가 시급함.", 6),

        H("적용 대상"),
        B("경기도 스마트공장 보급 확산 사업 대상 중소 전자부품 제조사의 QC 부서 및 생산 현장", 2),
        B("→ 공정 중간 단계에서 불량 징후 사전 감지 + 현장 작업자 즉각 조치 지원", 6),

        H("문제의 중요성"),
        BL("평균 불량률 3~8%, 불량 1건당 원자재·재작업 비용 수십만~수백만 원 직접 손실"),
        BL("분석 인력 부재 → 원인 분석 평균 2~3일 소요 → 동일 원인 반복 불량"),
        BL("공정 변수 간 복합 상호작용 → 경험만으로 근본 원인 파악 불가 → AI 필수 영역", 6),

        H("KPI / 기대 변화"),
        BL("불량 사전 감지율: 0%(사후 검사) → 70%+ (공정 중 사전 예측)"),
        BL("원인 분석 시간: 2~3일 → 실시간 (AI 자동 변수 기여도 제시)"),
        BL("연간 불량 손실 비용: 20%+ 절감 목표"),
    ])
    guide2.top = TB_TOP
    guide2.left = TB_LEFT
    guide2.width = TB_WIDTH
    guide2.height = TB_HEIGHT

print("Slide 2 (Problem Definition) - Done")


# ═══════════════════════════════════════════════
# SLIDE 3: 제안 솔루션 개요
# ═══════════════════════════════════════════════
slide3 = prs.slides[2]
guide3 = find_guide_textbox(slide3)
if guide3:
    clear_and_set(guide3.text_frame, [
        H("QualityLens — AI 품질 예측·원인 분석 통합 플랫폼"),
        B("공정 센서 데이터를 AI가 실시간 분석하여 불량 가능성을 사전 감지하고, 설명 가능 AI(XAI)로 '어떤 변수가 불량을 유발하는지' 현장 작업자 눈높이에서 시각 제시", 6),

        H("핵심 프레임워크 — Predict → Explain → Act"),
        BLB("Predict", "XGBoost 분류 모델 → 양품/불량 실시간 예측, MCC 기반 최적화로 클래스 불균형 극복"),
        BLB("Explain", "SHAP(TreeExplainer) → 예측 근거를 공정 변수별 기여도로 분해"),
        BLB("Act", "SHAP 기여도 + 양품 통계 기준(Threshold) 결합 → Rule-based 조치 가이드 즉시 제시", 6),

        H("플랫폼 목적"),
        B("데이터 분석 전문가 없이 현장 QC 담당자가 직접 활용하는 '셀프서비스형 AI 품질 관리 도구'", 2),
        B("→ 중소기업이 별도 데이터팀 없이 예방적 품질 관리 즉시 도입 가능", 6),

        H("전체 구성"),
        BLB("데이터 수집", "센서 파이프라인 (MVP: 오픈 데이터셋 기반 스트리밍 시뮬레이터)"),
        BLB("AI 엔진", "XGBoost + SHAP + Threshold Engine"),
        BLB("플랫폼 UI", "Streamlit 멀티페이지 대시보드 (모니터링·예측·분석·조치·이력)"),
        BLB("조치 지원", "알림 + 원클릭 수용 (이상 감지 시 권장 조치 자동 제시)"),
    ])
    guide3.top = TB_TOP
    guide3.left = TB_LEFT
    guide3.width = TB_WIDTH
    guide3.height = TB_HEIGHT

print("Slide 3 (Solution Overview) - Done")


# ═══════════════════════════════════════════════
# SLIDE 4: 주요 기능 정의
# ═══════════════════════════════════════════════
slide4 = prs.slides[3]
guide4 = find_guide_textbox(slide4)
if guide4:
    clear_and_set(guide4.text_frame, [
        H("기능 1 — 실시간 공정 모니터링 & 불량 예측 (Predict)"),
        B("센서 데이터 실시간 수집·시각화, XGBoost 모델이 배치별 양품/불량 확률 산출"),
        B("정상/경고/위험 3단계 컬러 코딩 → 비전문가도 즉시 이상 인지. MCC 기반 모델 평가 적용", 6),

        H("기능 2 — XAI 기반 불량 원인 분석 (Explain)"),
        B("SHAP으로 예측 근거를 공정 변수별 기여도로 분해"),
        B("글로벌 피처 중요도(전체 트렌드) + 개별 Waterfall(건별 원인) 동시 제공", 6),

        H("기능 3 — 조치 가이드 & 이력 관리 (Act)"),
        B("SHAP 지목 이상 변수 + 양품 통계 정상 범위(Threshold) 비교·분석"),
        B("→ 구체적 조치 방향 팝업 제시 (예: '온도 165~175°C로 하향 조정 요망')"),
        B("→ 원클릭 수용 버튼으로 조치 이력 자동 기록", 6),

        H("기능 간 연계 — 순환 루프"),
        B("예측 결과 → 원인 분석 입력 → 조치 가이드 근거 → 피드백 → 예측 모델 재학습"),
        B("종합 대시보드(P1)가 전체 흐름의 허브 역할", 6),

        H("본선 MVP"),
        B("P1(대시보드) + P2(예측) + P3(분석): 완전 구현"),
        B("P4(조치 가이드): 핵심 불량 시나리오 데모 구현 → End-to-End 사용자 경험 시연"),
    ])
    guide4.top = TB_TOP
    guide4.left = TB_LEFT
    guide4.width = TB_WIDTH
    guide4.height = TB_HEIGHT

print("Slide 4 (Key Features) - Done")


# ═══════════════════════════════════════════════
# SLIDE 5: 데이터 및 기술 활용 계획
# ═══════════════════════════════════════════════
slide5 = prs.slides[4]
guide5 = find_guide_textbox(slide5)
if guide5:
    clear_and_set(guide5.text_frame, [
        H("활용 데이터"),
        BLB("주 데이터셋", "UCI SECOM — 실제 반도체 제조 라인 (1,567샘플 × 591센서, 불량 104건/6.6%)"),
        BLB("보조 데이터셋", "Kaggle Smart Manufacturing — 직관적 공정 파라미터(온도/압력/진동/유량) 보강용", 6),

        H("데이터 처리"),
        BLB("결측치", "센서별 결측률 분석 → 고결측 피처 제거 + KNN Imputation"),
        BLB("피처 선택", "분산 필터링 + 상관관계 중복 제거 → 상위 50~80개 핵심 피처"),
        BLB("데이터 누수 방지", "Stratified Split → Train에만 SMOTE·정규화 적용 (엄격한 Pipeline 구축)", 6),

        H("AI/분석 기술"),
        BLB("XGBoost", "비선형 상호작용 학습, 빠른 추론, GPU 불필요"),
        BLB("SHAP (TreeExplainer)", "개별 예측의 피처별 기여도 실시간 분해"),
        BLB("Threshold Engine", "양품 통계 기준으로 조치 방향 및 권장 범위 자동 생성"),
        BLB("평가 지표", "MCC(주) + Precision-Recall, F1-Score(보조)", 4),

        SH("기술 스택: Streamlit | XGBoost | SHAP | Plotly | Pandas | scikit-learn | imbalanced-learn", 6),

        H("기술적 제약 → 해결 전략"),
        BL("컴퓨팅 미제공 → Streamlit + XGBoost 로컬 구동 (GPU 불필요)"),
        BL("실시간 연동 불가 → 오픈 데이터 기반 스트리밍 시뮬레이터"),
        BL("데이터 불균형(6.6%) → SMOTE + MCC + Stratified K-Fold"),
    ])
    guide5.top = TB_TOP
    guide5.left = TB_LEFT
    guide5.width = TB_WIDTH
    guide5.height = TB_HEIGHT

print("Slide 5 (Data & Tech Plan) - Done")


# ═══════════════════════════════════════════════
# SLIDE 6: 사용자 시나리오/유즈케이스
# ═══════════════════════════════════════════════
slide6 = prs.slides[5]
guide6 = find_guide_textbox(slide6)
if guide6:
    clear_and_set(guide6.text_frame, [
        H("주요 사용자"),
        BLB("페르소나 A — 김 과장 (QC)", "공정 경력 8년, 데이터 분석 무경험. 매일 P1~P4 전체 사용"),
        BLB("페르소나 B — 박 팀장 (생산)", "주간 리포트 검토, 의사결정권자. 주 1~2회 P1, P5 활용", 6),

        H("대표 시나리오: 식각 공정 온도 이상 → 불량 급증 대응"),
        SP(2),
        SH("Step 1 (08:30, P1 대시보드)"),
        B("출근 후 접속 → '금일 불량률 12.3% (+6.1%p)' 경고 + '배치 #B-2247~2251 연속 불량' 알림 확인", 4),

        SH("Step 2 (08:35, P2 실시간 예측)"),
        B("진행 중 배치 #B-2252 → AI 불량 확률 78% 산출, 온도 게이지 빨간색(위험) 표시", 4),

        SH("Step 3 (08:37, P3 원인 분석)"),
        B("SHAP Waterfall: '식각 온도(기여도 +0.34)' 1위 → Dependence Plot에서 180°C 이상 불량 급등 확인", 4),

        SH("Step 4 (08:40, P4 조치 가이드)"),
        B("팝업: '식각 온도 → 165~175°C 하향 조정 요망' → '수용' 버튼 클릭 → 이력 자동 기록", 4),

        SH("Step 5 (09:00, P1 확인)"),
        B("조치 후 #B-2253부터 불량률 2.1% 정상 복귀. 총 대응 시간: 약 30분 (기존 2~3일 대비 95%+ 단축)", 6),

        H("보조 시나리오: 주간 리포트"),
        B("박 팀장 → P5에서 주간 불량 추이, 반복 원인 Top 5 확인 → '식각 온도 3주 연속 주요 원인' → 설비 점검 앞당기기 의사결정"),
    ])
    guide6.top = TB_TOP
    guide6.left = TB_LEFT
    guide6.width = TB_WIDTH
    guide6.height = TB_HEIGHT

print("Slide 6 (User Case) - Done")


# ═══════════════════════════════════════════════
# SLIDE 7: MVP 구현 범위
# ═══════════════════════════════════════════════
slide7 = prs.slides[6]
guide7 = find_guide_textbox(slide7)
if guide7:
    clear_and_set(guide7.text_frame, [
        H("본선 구현 범위"),
        B("End-to-End 사용자 경험 완결 최소 범위. 원칙: '예측→조치까지 끊기지 않는 흐름'", 4),
        BLB("[필수] P1 대시보드", "KPI 카드 + 알림 패널 + 불량률 추이 — 완전 구현"),
        BLB("[필수] P2 실시간 예측", "센서 게이지 + 배치별 확률 + 시뮬레이터 — 완전 구현"),
        BLB("[필수] P3 원인 분석", "SHAP Summary + Waterfall + Dependence Plot — 완전 구현"),
        BLB("[핵심 데모] P4 조치 가이드", "핵심 시나리오 2~3개 + 원클릭 수용 — 시나리오 기반"),
        BLB("[확장] P5 이력 관리", "조치 이력 테이블 — 시간 여유 시 구현", 6),

        H("데모 시나리오 (약 5분)"),
        B("정상 스트리밍 → 이상 주입 → P1 경고 → P2 불량 확률 급등 → P3 원인 지목 → P4 조치 수용 → P1 정상 복귀"),
        B("→ Predict → Explain → Act 전체 루프를 한 번의 시연으로 완결", 4),
        ACC("※ Fail-Safe: 인프라 장애 대비, MVP 정상 구동 시연 영상 사전 준비", 6),

        H("사전 개발 전략 (예선 통과 후 ~ 본선 전 9일)"),
        BL("~5/18: 데이터 전처리 파이프라인 + XGBoost 학습·저장"),
        BL("~5/19: SHAP Explainer 사전 계산·캐싱"),
        BL("~5/20: Streamlit 페이지 골격 (레이아웃 + 네비게이션)"),
        BL("~5/21: 시뮬레이터 데이터 세트 준비"),
        B("→ 본선 당일: '조립 + 연결 + 시각적 완성도'에 집중", 2),
    ])
    guide7.top = TB_TOP
    guide7.left = TB_LEFT
    guide7.width = TB_WIDTH
    guide7.height = TB_HEIGHT

print("Slide 7 (MVP Scope) - Done")


# ═══════════════════════════════════════════════
# SLIDE 8: 기대 효과 및 향후 확장성
# ═══════════════════════════════════════════════
slide8 = prs.slides[7]
guide8 = find_guide_textbox(slide8)
if guide8:
    clear_and_set(guide8.text_frame, [
        H("정량적 기대 효과"),
        BL("불량 감지: 최종 검사(사후) → 공정 중간(사전) — 70%+ 사전 감지"),
        BL("원인 분석: 2~3일(수기) → 실시간 자동 — 99% 단축"),
        BL("대응 시간: 수일 → 30분 이내 — 95%+ 단축"),
        BL("손실 비용: 연간 20%+ 절감 (중소기업 기준 수천만 원)", 6),

        H("정성적 기대 효과"),
        BL("데이터 기반 품질 문화 정착: 경험·직감 → AI 근거 기반 의사결정 전환"),
        BL("현장 역량 강화: AI가 대체가 아닌 보강(Augmentation), SHAP으로 공정 학습"),
        BL("중소기업 진입 장벽 해소: 데이터팀 없이 Python 오픈소스로 즉시 도입 가능", 6),

        H("현장 적용 가능성"),
        B("MVP → 현장 전환 비용이 낮음. 데이터 소스만 교체(오픈 데이터 → MES/센서 API)하면 AI 파이프라인·UI 그대로 유지", 6),

        H("향후 고도화 로드맵"),
        BLB("Phase 1 (1~3개월)", "다중 공정 지원 + 자동 재학습(Auto-Retrain)"),
        BLB("Phase 2 (3~6개월)", "설비 고장 예지 모듈 → 품질+설비 통합 관리"),
        BLB("Phase 3 (6~12개월)", "생산 계획 최적화 + 안전 관리 → 중소기업형 통합 스마트 공장 플랫폼"),
    ])
    guide8.top = TB_TOP
    guide8.left = TB_LEFT
    guide8.width = TB_WIDTH
    guide8.height = TB_HEIGHT

print("Slide 8 (Expected Impact) - Done")


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
