#!/usr/bin/env python3
"""Build QualityLens 기획서 PPTX from template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = "/sessions/clever-nifty-volta/mnt/Dacon/smart-factory-hackathon/2026-산출물-양식.pptx"
OUTPUT = "/sessions/clever-nifty-volta/mnt/Dacon/smart-factory-hackathon/QualityLens_기획서.pptx"

prs = Presentation(TEMPLATE)

# ═══════════════════════════════════════════════
# Helper: clear & set text in a text frame
# ═══════════════════════════════════════════════
def clear_and_set(text_frame, paragraphs_data, font_size=11, font_name="Malgun Gothic"):
    """
    paragraphs_data: list of dicts
      { "text": str, "bold": bool, "color": RGBColor or None, "size": Pt or None, "bullet": bool }
    """
    # Clear existing paragraphs
    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
        p = text_frame.paragraphs[i]._p
        p.getparent().remove(p)
    
    first = True
    for pd in paragraphs_data:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        p.text = pd.get("text", "")
        p.font.size = pd.get("size", Pt(font_size))
        p.font.name = font_name
        p.font.bold = pd.get("bold", False)
        p.font.color.rgb = pd.get("color", RGBColor(0x20, 0x21, 0x24))
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        
        if pd.get("indent", 0) > 0:
            p.level = pd.get("indent", 0)
        
        # Line spacing
        p.space_after = Pt(pd.get("space_after", 4))


def find_guide_textbox(slide):
    """Find the [작성방법] guide textbox in a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            if "작성방법" in full_text:
                return shape
    return None

def find_content_rect(slide):
    """Find the bordered rectangle content area."""
    for shape in slide.shapes:
        if hasattr(shape, 'shape_id'):
            name = shape.name or ""
            if "직사각형" in name:
                return shape
    return None


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_table:
        table = shape.table
        # Row 0: 팀명
        table.cell(0, 1).text = "QualityLens"
        for p in table.cell(0, 1).text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = "Malgun Gothic"
            p.font.bold = True
        # Row 1: 프로젝트명
        table.cell(1, 1).text = "QualityLens — AI 기반 스마트 품질 예측·원인 분석 통합 플랫폼"
        for p in table.cell(1, 1).text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Malgun Gothic"

print("Slide 1 (Title) - Done")


# ═══════════════════════════════════════════════
# SLIDE 2: 문제 정의 (Problem Definition)
# ═══════════════════════════════════════════════
slide2 = prs.slides[1]
guide2 = find_guide_textbox(slide2)
if guide2:
    tf = guide2.text_frame
    clear_and_set(tf, [
        {"text": "해결하고자 하는 제조 현장의 문제", "bold": True, "size": Pt(14), "color": RGBColor(0x20, 0x21, 0x24), "space_after": 2},
        {"text": "경기도 내 중소 전자부품·반도체 부품 제조사들은 사출, 도금, 식각, 열처리 등 복합 공정에서 발생하는 품질 불량에 대해 최종 검사 이후에야 불량을 확인하는 사후 대응 방식에 의존하고 있습니다. 각 공정 단계에서 수집되는 수백 개의 센서 데이터가 존재하지만, 전문 데이터 분석 인력이 부족한 중소기업 특성상 이 데이터를 품질 관리에 활용하지 못하고 있습니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "적용 대상 및 영역", "bold": True, "size": Pt(14), "color": RGBColor(0x20, 0x21, 0x24), "space_after": 2},
        {"text": "경기도 스마트공장 보급 확산 사업 대상인 중소 전자부품 제조사의 품질관리(QC) 부서 및 생산 현장. 특히 공정 중간 단계에서 불량 징후를 사전 감지하고, 현장 작업자가 즉각적인 공정 조건 조정을 할 수 있도록 지원하는 '예방적 품질 관리 플랫폼'을 목표로 합니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "문제의 중요성", "bold": True, "size": Pt(14), "color": RGBColor(0x20, 0x21, 0x24), "space_after": 2},
        {"text": "• 중소 제조기업의 평균 불량률은 3~8%이며, 불량 1건당 원자재·설비 가동 시간·재작업 비용 포함 수십만~수백만 원의 직접 손실 발생", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 전문 분석 인력 부재로 불량 원인 분석에 평균 2~3일 소요 → 같은 원인으로 반복 불량 발생", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 공정 변수 간 복합적 상호작용이 불량을 유발하므로 경험만으로는 근본 원인 파악이 어려움 → AI가 반드시 필요한 영역", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "KPI / 기대 변화", "bold": True, "size": Pt(14), "color": RGBColor(0x20, 0x21, 0x24), "space_after": 2},
        {"text": "• 불량 사전 감지율: 0%(사후 검사) → 70%+ 공정 중간에서 사전 예측", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 불량 원인 분석 시간: 2~3일 → 실시간 (AI가 핵심 공정 변수와 기여도를 자동 제시)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 연간 불량 관련 손실 비용: 20% 이상 절감 목표", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    # Resize the textbox to fill the content area
    guide2.top = Emu(1600000)
    guide2.left = Emu(700000)
    guide2.width = Emu(10800000)
    guide2.height = Emu(4800000)

print("Slide 2 (Problem Definition) - Done")


# ═══════════════════════════════════════════════
# SLIDE 3: 제안 솔루션 개요 (Solution Overview)
# ═══════════════════════════════════════════════
slide3 = prs.slides[2]
guide3 = find_guide_textbox(slide3)
if guide3:
    tf = guide3.text_frame
    clear_and_set(tf, [
        {"text": "제안하는 솔루션의 개요", "bold": True, "size": Pt(14), "space_after": 2},
        {"text": "본 프로젝트는 'QualityLens' — 공정 센서 데이터 기반의 AI 품질 예측·원인 분석 통합 플랫폼을 제안합니다. 공정 진행 중 실시간으로 수집되는 센서 데이터를 AI 모델이 분석하여 불량 가능성이 높은 제품을 사전에 감지하고, 설명 가능 AI(XAI)를 통해 '어떤 공정 변수가 불량을 유발하고 있는지'를 현장 작업자 눈높이에서 시각적으로 제시합니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "문제 해결을 위한 핵심 아이디어 — Predict → Explain → Act", "bold": True, "size": Pt(14), "space_after": 2},
        {"text": "• Predict(예측): XGBoost 기반 분류 모델이 공정 센서 데이터를 입력받아 양품/불량을 실시간으로 예측. 제조 데이터의 극심한 클래스 불균형을 극복하기 위해 MCC 지표 기반 최적화 적용", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• Explain(설명): SHAP(SHapley Additive exPlanations)이 예측 결과의 근거를 공정 변수별 기여도로 분해하여 구체적으로 제시", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• Act(조치): SHAP 기여도 분석과 정상 양품 데이터의 통계적 기준(Threshold)을 결합한 Rule-based 조치 가이드를 대시보드에 즉각 표시", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "플랫폼의 목적 및 제공 가치", "bold": True, "size": Pt(14), "space_after": 2},
        {"text": "데이터 분석 전문가 없이도 현장 품질관리 담당자가 직접 활용할 수 있는 '셀프서비스형 AI 품질 관리 도구'를 제공합니다. 중소기업이 별도의 데이터팀 구축 없이 예방적 품질 관리를 즉시 도입할 수 있도록 하는 것이 목표입니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "전체 구성 요약", "bold": True, "size": Pt(14), "space_after": 2},
        {"text": "• 데이터 수집 — 센서 데이터 파이프라인 (MVP: 오픈 데이터셋 기반 실시간 스트리밍 시뮬레이터)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• AI 엔진 — XGBoost + SHAP + Threshold Engine (불량 예측 + 원인 분석 + 조치 근거)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 플랫폼 UI — Streamlit 멀티페이지 대시보드 (모니터링·예측·분석·조치·이력 관리)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 조치 지원 — 알림 + 원클릭 조치 가이드 (이상 감지 시 작업자 알림 및 권장 조치 제시)", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    guide3.top = Emu(1600000)
    guide3.left = Emu(700000)
    guide3.width = Emu(10800000)
    guide3.height = Emu(4800000)

print("Slide 3 (Solution Overview) - Done")


# ═══════════════════════════════════════════════
# SLIDE 4: 주요 기능 정의 (Key Features)
# ═══════════════════════════════════════════════
slide4 = prs.slides[3]
guide4 = find_guide_textbox(slide4)
if guide4:
    tf = guide4.text_frame
    clear_and_set(tf, [
        {"text": "기능 1 — 실시간 공정 모니터링 & 불량 예측 (Predict)", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "공정 센서 데이터를 실시간으로 수집·시각화하고, XGBoost 모델이 배치별 양품/불량 확률을 산출합니다. 핵심 센서의 정상/경고/위험 상태를 3단계 컬러 코딩으로 표시하여, 비전문가도 즉시 공정 이상을 인지할 수 있습니다. 클래스 불균형 환경에서의 신뢰성 확보를 위해 MCC 기반 모델 평가를 적용합니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "기능 2 — XAI 기반 불량 원인 분석 (Explain)", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "SHAP을 활용하여 모델의 예측 근거를 공정 변수별 기여도로 분해합니다. 글로벌 피처 중요도(전체 트렌드)와 개별 Waterfall(건별 원인)을 모두 제공하여, '전체적으로 어떤 변수가 중요한지'와 '이 배치가 왜 불량인지'를 동시에 파악할 수 있습니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "기능 3 — 조치 가이드 & 이력 관리 (Act)", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "SHAP이 지목한 핵심 이상 변수와 과거 양품 데이터의 통계적 정상 범위(Threshold)를 비교·분석하여, 현장 작업자에게 구체적인 조치 방향(예: '온도 센서 비정상 상승 감지 → 정상 권장 범위 165~175°C로 하향 조정 요망')을 대시보드에 즉각 팝업 형태로 제시합니다. 원클릭 수용 버튼으로 조치 이력이 자동 기록됩니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "기능 간 연계 방식", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "세 기능은 독립적이지 않고 순환 루프를 구성합니다 — 예측 결과가 원인 분석의 입력이 되고, 분석 결과가 조치 가이드의 근거가 되며, 조치 결과가 다시 예측 모델의 학습 데이터로 피드백됩니다. 종합 대시보드(P1)가 이 전체 흐름의 허브 역할을 합니다.", "bold": False, "size": Pt(10), "space_after": 8},
        {"text": "본선 MVP 핵심 기능", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "P1(종합 대시보드), P2(실시간 예측), P3(원인 분석)의 완전한 파이프라인 구현을 필수로 진행합니다. P4(조치 가이드)는 핵심 불량 시나리오를 가정한 데모 화면을 구성하여, 예측부터 조치까지의 End-to-End 사용자 경험을 완벽하게 시연할 계획입니다.", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    guide4.top = Emu(1600000)
    guide4.left = Emu(700000)
    guide4.width = Emu(10800000)
    guide4.height = Emu(4800000)

print("Slide 4 (Key Features) - Done")


# ═══════════════════════════════════════════════
# SLIDE 5: 데이터 및 기술 활용 계획
# ═══════════════════════════════════════════════
slide5 = prs.slides[4]
guide5 = find_guide_textbox(slide5)
if guide5:
    tf = guide5.text_frame
    clear_and_set(tf, [
        {"text": "활용할 데이터 종류 및 특성", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 주 데이터셋: UCI SECOM — 실제 반도체 제조 라인 데이터 (1,567샘플 × 591센서, 불량 104건/6.6%)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 보조 데이터셋: Kaggle Smart Manufacturing Process Data — 직관적 공정 파라미터(온도/압력/진동/유량) 보강용", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "데이터 처리 방식", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 결측치 처리: 센서별 결측률 분석 → 고결측 피처 제거 + KNN Imputation", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 피처 선택: 분산 기반 필터링 + 상관관계 기반 중복 제거 → 상위 50~80개 핵심 피처 추출", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 클래스 밸런싱 & 데이터 누수 방지: Stratified Split 이후 Train 데이터에만 SMOTE·정규화 적용하는 엄격한 Pipeline 구축으로 Data Leakage 차단", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "적용할 AI/분석 기술", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• XGBoost 분류 모델: 센서 데이터의 비선형 상호작용 학습, 빠른 추론 속도, GPU 불필요", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• SHAP (TreeExplainer): XGBoost 전용 고속 SHAP 연산으로 개별 예측의 피처별 기여도 실시간 분해", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• Threshold Engine: 양품 데이터의 센서별 통계 기준으로 조치 방향과 권장 범위 자동 생성", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 평가 지표: MCC(Matthews Correlation Coefficient) 주요 지표 + Precision-Recall 커브, F1-Score 보조", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "기술 스택: Streamlit | XGBoost | SHAP | Plotly | Pandas | scikit-learn | imbalanced-learn", "bold": True, "size": Pt(11), "space_after": 6},
        {"text": "기술적 제약 및 해결 전략", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 컴퓨팅 리소스 미제공 → Streamlit + XGBoost는 로컬 노트북에서 충분히 구동 (GPU 불필요)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 실시간 센서 연동 불가 → 오픈 데이터셋 기반 스트리밍 시뮬레이터로 실시간 환경 재현", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 데이터 불균형 (6.6% 불량) → SMOTE + MCC 평가 + Stratified K-Fold 교차 검증", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    guide5.top = Emu(1600000)
    guide5.left = Emu(700000)
    guide5.width = Emu(10800000)
    guide5.height = Emu(4800000)

print("Slide 5 (Data & Tech Plan) - Done")


# ═══════════════════════════════════════════════
# SLIDE 6: 사용자 시나리오/유즈케이스
# ═══════════════════════════════════════════════
slide6 = prs.slides[5]
guide6 = find_guide_textbox(slide6)
if guide6:
    tf = guide6.text_frame
    clear_and_set(tf, [
        {"text": "주요 사용자 정의", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 페르소나 A — 김 과장 (QC 담당자): 공정 경력 8년, 데이터 분석 경험 없음, 엑셀 수준 IT 역량. 매일 P1~P4 전체 사용", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 페르소나 B — 박 팀장 (생산팀장): 주간 품질 리포트 검토, 의사결정권자. 주 1~2회 P1, P5 활용", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "대표 시나리오: 식각 공정 온도 이상 → 불량 급증 대응", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "Step 1 (08:30, P1 대시보드) — 김 과장 출근 후 QualityLens 접속. '금일 불량률 12.3% (+6.1%p)' 경고 확인. 알림 패널에서 '식각 공정 배치 #B-2247~#B-2251 연속 불량 감지' 알림 확인", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "Step 2 (08:35, P2 실시간 예측) — 진행 중 배치 #B-2252 센서 게이지 실시간 확인. AI 모델 '불량 확률 78%' 산출, 온도 게이지 빨간색(위험) 표시", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "Step 3 (08:37, P3 원인 분석) — SHAP Waterfall에서 '식각 온도(기여도 +0.34)' 1위 원인 지목. Dependence Plot에서 180°C 이상 구간 불량 급등 패턴 시각적 확인", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "Step 4 (08:40, P4 조치 가이드) — 팝업: '식각 온도 비정상 상승 → 정상 권장 범위 165~175°C로 하향 조정 요망'. 김 과장이 '수용' 버튼 클릭 → 조치 이력 P5에 자동 기록", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "Step 5 (09:00~, P1 확인) — 조치 후 #B-2253부터 불량률 2.1%로 정상 복귀 확인. 총 대응 시간: 약 30분 (기존 2~3일 → 95%+ 단축)", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "보조 시나리오: 박 팀장의 주간 리포트", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "매주 월요일 P5(이력 관리)에서 주간 불량률 추이, AI 감지 정확도, 반복 원인 Top 5를 확인. '식각 온도 이상이 3주 연속 주요 원인' 패턴 발견 → 설비 점검 일정 앞당기는 의사결정", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    guide6.top = Emu(1600000)
    guide6.left = Emu(700000)
    guide6.width = Emu(10800000)
    guide6.height = Emu(4800000)

print("Slide 6 (User Case) - Done")


# ═══════════════════════════════════════════════
# SLIDE 7: MVP 구현 범위
# ═══════════════════════════════════════════════
slide7 = prs.slides[6]
guide7 = find_guide_textbox(slide7)
if guide7:
    tf = guide7.text_frame
    clear_and_set(tf, [
        {"text": "본선에서 구현할 범위", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "End-to-End 사용자 경험이 완결되는 최소 범위를 MVP로 설정합니다. 핵심 원칙: '예측부터 조치까지 끊기지 않는 흐름'", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "• [필수] P1 종합 대시보드: KPI 카드 + 알림 패널 + 불량률 추이 차트 — 완전 구현", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• [필수] P2 실시간 예측: 센서 게이지 + 배치별 불량 확률 + 시뮬레이터 스트리밍 — 완전 구현", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• [필수] P3 원인 분석: SHAP Summary + Waterfall + Dependence Plot — 완전 구현", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• [핵심 데모] P4 조치 가이드: 핵심 불량 시나리오 2~3개 + 원클릭 수용 버튼 — 시나리오 기반 구현", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• [확장] P5 이력 관리: 조치 이력 테이블 — 시간 여유 시 구현", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "데모 시나리오 (약 5분)", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "시뮬레이터 정상 스트리밍 → 이상 데이터 주입 → P1 경고 → P2 불량 확률 급등 → P3 SHAP 원인 지목 → P4 조치 가이드 확인 후 원클릭 수용 → P1 불량률 정상 복귀 확인. Predict → Explain → Act 전체 루프를 한 번의 시연으로 완결.", "bold": False, "size": Pt(10), "space_after": 4},
        {"text": "※ Fail-Safe: 네트워크 불안정 등 인프라 장애 대비, 사전 개발된 MVP 정상 구동 시연 영상을 별도 준비하여 평가 안정성 확보", "bold": True, "size": Pt(10), "color": RGBColor(0xC6, 0x28, 0x28), "space_after": 6},
        {"text": "사전 개발 전략 (예선 통과 후 ~ 본선 전 9일간)", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• ~5/18: 데이터 전처리 파이프라인 + XGBoost 모델 학습·저장 완료", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• ~5/19: SHAP Explainer 사전 계산·캐싱 완료", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• ~5/20: Streamlit 페이지 골격(레이아웃+네비게이션) 완성", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• ~5/21: 시뮬레이터(정상/이상 시나리오) 데이터 세트 준비", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "→ 본선 당일에는 '조립 + 연결 + 시각적 완성도 높이기'에 집중", "bold": True, "size": Pt(10), "space_after": 2},
    ])
    guide7.top = Emu(1600000)
    guide7.left = Emu(700000)
    guide7.width = Emu(10800000)
    guide7.height = Emu(4800000)

print("Slide 7 (MVP Scope) - Done")


# ═══════════════════════════════════════════════
# SLIDE 8: 기대 효과 및 향후 확장성
# ═══════════════════════════════════════════════
slide8 = prs.slides[7]
guide8 = find_guide_textbox(slide8)
if guide8:
    tf = guide8.text_frame
    clear_and_set(tf, [
        {"text": "정량적 기대 효과", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 불량 감지 시점: 최종 검사(사후) → 공정 중간(사전 예측) — 70%+ 사전 감지", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 불량 원인 분석 시간: 평균 2~3일(수기) → 실시간 자동 분석 — 99% 단축", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 불량 대응 총 소요 시간: 수일 → 30분 이내 — 95%+ 단축", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 연간 불량 관련 손실 비용: 20% 이상 절감 (중소기업 기준 수천만 원)", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "정성적 기대 효과", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• 데이터 기반 품질 문화 정착: 경험·직감 의존 → AI 근거 기반 의사결정 전환. 조치 이력 자동 축적으로 조직 지식 체계화", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 현장 작업자 역량 강화: AI가 작업자를 대체하는 것이 아니라 보강(Augmentation). SHAP 시각화를 통해 공정 변수 간 관계를 자연스럽게 학습", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• 중소기업 스마트 팩토리 진입 장벽 해소: 별도 데이터팀 없이 Python 기반 오픈소스 스택으로 즉시 도입 가능. 경기도 스마트공장 보급 확산 사업 취지와 정확히 부합", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "현장 적용 가능성", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "MVP에서 현장 도입으로의 전환 비용이 낮습니다. Streamlit이 Python 생태계 위에서 동작하므로 데이터 소스만 교체(오픈 데이터셋 → MES/센서 API)하면 AI 파이프라인과 UI가 그대로 유지됩니다.", "bold": False, "size": Pt(10), "space_after": 6},
        {"text": "향후 고도화 로드맵", "bold": True, "size": Pt(13), "space_after": 2},
        {"text": "• Phase 1 (도입 후 1~3개월): 다중 공정 지원 + 자동 재학습(Auto-Retrain)", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• Phase 2 (3~6개월): 설비 고장 예지 모듈 추가 → 품질+설비 통합 관리", "bold": False, "size": Pt(10), "space_after": 2},
        {"text": "• Phase 3 (6~12개월): 생산 계획 최적화 + 안전 관리 → 중소기업형 통합 스마트 공장 플랫폼으로 확장", "bold": False, "size": Pt(10), "space_after": 2},
    ])
    guide8.top = Emu(1600000)
    guide8.left = Emu(700000)
    guide8.width = Emu(10800000)
    guide8.height = Emu(4800000)

print("Slide 8 (Expected Impact) - Done")

# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
